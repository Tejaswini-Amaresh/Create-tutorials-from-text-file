from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from nltk.tokenize import word_tokenize, sent_tokenize

import nltk

import re

def processing(text,filename):


  '''CREATING FREQUENCY TABLE'''
  stopWords = set(stopwords.words("english"))

  #print(stopWords)

  import nltk
  nltk.download('punkt')
  words = word_tokenize(text)

  #print(words)

  ps = PorterStemmer()

  freqTable = dict()
  for word in words:
      word = ps.stem(word)
      if word in stopWords:
              continue
      if word in freqTable:
              freqTable[word] += 1
      else:
              freqTable[word] = 1

  #freqTable

  '''SCORING SENTENCES'''
  sentenceValue = dict()

  sentences = text.split('.')

  for i in range(len(sentences)):
    sentences[i]=sentences[i]+"."

  #sentences

  #len(sentences)

  for sentence in sentences:
          word_count_in_sentence = (len(word_tokenize(sentence)))
          word_count_in_sentence_except_stop_words = 0
          for wordValue in freqTable:
              if wordValue in sentence.lower():
                  word_count_in_sentence_except_stop_words += 1
                  if sentence[:10] in sentenceValue:
                      sentenceValue[sentence[:10]] += freqTable[wordValue]
                  else:
                      sentenceValue[sentence[:10]] = freqTable[wordValue]

  if sentence[:10] in sentenceValue:
          sentenceValue[sentence[:10]] = sentenceValue[sentence[:10]]

  #sentenceValue

  '''FINDING AVERAGE SCORE'''
  sumValues = 0
  for entry in sentenceValue:
      sumValues += sentenceValue[entry]

      # Average value of a sentence from original text
      average = (sumValues / len(sentenceValue))

  #average

  sentence_count = 0
  summary = ''
  mapping={}
  for i in range(len(sentences)):
      sentence=sentences[i]
      if sentence[:10] in sentenceValue and sentenceValue[sentence[:10]] >= (average):
              summary += " " + sentence
              sentence_count += 1
              mapping[sentence_count-1]=i
  
  print(mapping)

  #summary



  from gtts import gTTS 
  mytext = summary
  language = 'en'
  
  myobj = gTTS(text=text, lang=language, slow=False) 

  myobj.save("voiceover.mp3")

  summary = summary.split('.')
  summary.pop()

  for i in range(len(summary)):
    summary[i]=summary[i]+"."

  #summary

  from pptx import Presentation,util

  prs = Presentation()
  title_slide_layout = prs.slide_layouts[1]
  '''slide = prs.slides.add_slide(title_slide_layout)
  title = slide.shapes.title
  subtitle = slide.placeholders[1]
  left = top = width = height = util.Inches(1.0)
  shapes = slide.shapes
  body_shape = shapes.placeholders[1]
  title_shape = shapes.title'''
  #Used to add bullet points
  #tf = body_shape.text_frame
  for i in range(0,len(summary),3):
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    left = top = width = height = util.Inches(1.0)
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    title_shape = shapes.title
    x=min(i+3,len(summary))
    for j in range(i,x):
      p = tf.add_paragraph()
      p.text = summary[j]
      p.level=0
    voice_text=' '.join(sentences[mapping[i]:mapping[x-1]+1])
    myobj = gTTS(text=voice_text, lang=language, slow=False) 
    myobj.save("voiceover.mp3")
    movie = slide.shapes.add_movie("voiceover.mp3", 
                    left , top , width , height, 
                              poster_frame_image=None, 
                              mime_type='video/unknown')

  prs.save('static/downloads/'+filename)

